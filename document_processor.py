"""
document_processor.py — Text Extraction and Chunking (Phase 2)

Extracts text from PDF, DOCX, PPTX, and TXT files.
Chunks text into overlapping segments for embedding models with 512-token limits.
"""

import os


def extract_pdf(file_path):
    """Extract all text from a PDF file using PyMuPDF."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        text_parts = []
        for page in doc:
            text_parts.append(page.get_text())
        doc.close()
        return "\n".join(text_parts).strip()
    except Exception as e:
        print(f"  Error extracting PDF '{file_path}': {e}")
        return ""


def extract_docx(file_path):
    """Extract all text from a Word DOCX file."""
    try:
        from docx import Document
        doc = Document(file_path)
        text_parts = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n".join(text_parts).strip()
    except Exception as e:
        print(f"  Error extracting DOCX '{file_path}': {e}")
        return ""


def extract_pptx(file_path):
    """Extract all text from a PowerPoint PPTX file."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            text_parts.append(text)
        return "\n".join(text_parts).strip()
    except Exception as e:
        print(f"  Error extracting PPTX '{file_path}': {e}")
        return ""


def extract_txt(file_path):
    """Extract text from a plain text file with encoding fallback."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read().strip()
    except UnicodeDecodeError:
        try:
            with open(file_path, "r", encoding="latin-1") as f:
                return f.read().strip()
        except Exception as e:
            print(f"  Error extracting TXT '{file_path}': {e}")
            return ""
    except Exception as e:
        print(f"  Error extracting TXT '{file_path}': {e}")
        return ""


def extract_text(file_path):
    """
    Dispatcher: routes to the correct extraction function based on file extension.
    Returns the extracted text as a string (empty string on failure).
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == ".pdf":
        return extract_pdf(file_path)
    elif ext in (".doc", ".docx"):
        return extract_docx(file_path)
    elif ext in (".ppt", ".pptx"):
        return extract_pptx(file_path)
    elif ext in (".txt", ".md", ".csv", ".log"):
        return extract_txt(file_path)
    else:
        return ""


def chunk_text(text, chunk_size=500, overlap=50):
    """
    Split text into overlapping word-level chunks.
    
    Args:
        text: The full text to chunk.
        chunk_size: Number of words per chunk. 
        overlap: Number of overlapping words between consecutive chunks.
    
    Returns:
        List of chunk strings. Returns [text] if text is shorter than chunk_size.
    """
    if not text or not text.strip():
        return []
    
    words = text.split()
    
    # If the text is short enough, return it as a single chunk
    if len(words) <= chunk_size:
        return [text.strip()]
    
    chunks = []
    start = 0
    
    while start < len(words):
        end = start + chunk_size
        chunk_words = words[start:end]
        chunk = " ".join(chunk_words)
        
        if chunk.strip():
            chunks.append(chunk.strip())
        
        # Move forward by (chunk_size - overlap) to create overlap
        start += chunk_size - overlap
    
    return chunks


if __name__ == "__main__":
    # Quick self-test
    print("=== Document Processor Self-Test ===\n")
    
    sample_text = " ".join([f"word{i}" for i in range(1200)])
    chunks = chunk_text(sample_text, chunk_size=500, overlap=50)
    print(f"Input: {len(sample_text.split())} words")
    print(f"Output: {len(chunks)} chunks")
    for i, c in enumerate(chunks):
        print(f"  Chunk {i}: {len(c.split())} words | starts with '{' '.join(c.split()[:3])}...'")
